"""Multi-format document processor."""

import csv
import uuid
import io
import zipfile
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Any, List
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import anthropic

from src.models.document import ProcessingStatus, DocumentCategory
from src.config import settings
from src.utils.logging import get_logger

logger = get_logger(__name__)


class MultiFormatProcessor:
    """Process documents in multiple formats."""

    def __init__(self):
        """Initialize processor."""
        try:
            self.anthropic_client = anthropic.Anthropic(api_key=settings.anthropic_api_key)
        except Exception as e:
            logger.warning(f"Failed to initialize Anthropic client: {e}")
            self.anthropic_client = None

    def detect_file_type(self, file_path: Path) -> str:
        """
        Detect file type from extension and content.
        For files without extension, uses file signature (magic numbers).

        Args:
            file_path: Path to file

        Returns:
            str: File type (pdf, csv, xml, html, image, text, docx, xlsx, pptx, json, unknown)
        """
        suffix = file_path.suffix.lower()

        type_mapping = {
            '.pdf': 'pdf',
            '.csv': 'csv',
            '.xml': 'xml',
            '.xaf': 'xaf',
            '.html': 'html',
            '.htm': 'html',
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.gif': 'image',
            '.bmp': 'image',
            '.svg': 'image',
            '.webp': 'image',
            '.txt': 'text',
            '.md': 'text',
            '.markdown': 'text',
            '.rst': 'text',
            '.log': 'text',
            '.json': 'json',
            '.docx': 'docx',
            '.doc': 'doc',
            '.xlsx': 'xlsx',
            '.xls': 'xls',
            '.pptx': 'pptx',
            '.ppt': 'ppt',
            '.odt': 'odt',
            '.ods': 'ods',
            '.odp': 'odp',
            '.rtf': 'rtf'
        }

        # If extension is known, return it
        if suffix and suffix in type_mapping:
            return type_mapping[suffix]

        # For files without extension or unknown extension, check file signature
        try:
            with open(file_path, 'rb') as f:
                header = f.read(16)

            if not header:
                return 'unknown'

            # Check magic numbers
            # PNG
            if header.startswith(b'\x89PNG\r\n\x1a\n'):
                return 'image'
            # JPEG
            elif header.startswith(b'\xff\xd8\xff'):
                return 'image'
            # GIF
            elif header.startswith(b'GIF87a') or header.startswith(b'GIF89a'):
                return 'image'
            # PDF
            elif header.startswith(b'%PDF'):
                return 'pdf'
            # ZIP-based formats (DOCX, XLSX, PPTX)
            elif header.startswith(b'PK\x03\x04'):
                # Could be docx, xlsx, pptx, or just zip
                # Try to determine by reading more
                return 'docx'  # Default to docx, will be refined in extraction
            # JSON (starts with { or [)
            elif header.startswith(b'{') or header.startswith(b'['):
                return 'json'
            # XML
            elif header.startswith(b'<?xml') or header.startswith(b'<'):
                return 'xml'
            # HTML
            elif b'<html' in header.lower() or b'<!doctype html' in header.lower():
                return 'html'

        except Exception as e:
            logger.debug(f"Error detecting file type for {file_path.name}: {e}")

        return 'unknown'

    def extract_text_from_csv(self, file_path: Path) -> str:
        """
        Extract text content from CSV file.

        Args:
            file_path: Path to CSV file

        Returns:
            str: Formatted text representation of CSV data
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.DictReader(f)
                rows = list(reader)

                if not rows:
                    return "Empty CSV file"

                # Format as readable text
                text_parts = [f"CSV Data ({len(rows)} rows):\n"]

                # Add headers
                headers = rows[0].keys()
                text_parts.append(f"Columns: {', '.join(headers)}\n\n")

                # Add sample rows (first 10)
                for i, row in enumerate(rows[:10], 1):
                    text_parts.append(f"Row {i}:\n")
                    for key, value in row.items():
                        if value:  # Only include non-empty values
                            text_parts.append(f"  {key}: {value}\n")
                    text_parts.append("\n")

                if len(rows) > 10:
                    text_parts.append(f"... and {len(rows) - 10} more rows")

                return ''.join(text_parts)

        except Exception as e:
            logger.error(f"Error extracting text from CSV: {e}")
            return f"Error reading CSV: {str(e)}"

    def extract_text_from_xml(self, file_path: Path) -> str:
        """
        Extract text content from XML file.

        Args:
            file_path: Path to XML file

        Returns:
            str: Formatted text representation of XML data
        """
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            return self._render_xml_tree(root, f"XML Document: {root.tag}")

        except Exception as e:
            logger.error(f"Error extracting text from XML: {e}")
            return f"Error reading XML: {str(e)}"

    def extract_text_from_xaf(self, file_path: Path) -> str:
        """
        Extract text content from XAF (XML Auditfile Financieel).

        Supports plain XML .xaf files and zipped containers with XML members.
        """
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            return self._render_xml_tree(root, f"XAF Auditfile: {root.tag}")
        except ET.ParseError:
            if zipfile.is_zipfile(file_path):
                try:
                    with zipfile.ZipFile(file_path, "r") as archive:
                        xml_members = [
                            name for name in archive.namelist()
                            if name.lower().endswith((".xml", ".xaf"))
                        ]
                        if not xml_members:
                            return (
                                f"XAF archive: {file_path.name}\n"
                                "No XML members found in archive."
                            )

                        combined_parts = []
                        for member_name in xml_members:
                            try:
                                xml_bytes = archive.read(member_name)
                                root = ET.fromstring(xml_bytes)
                                combined_parts.append(
                                    self._render_xml_tree(root, f"XAF Member: {member_name}")
                                )
                            except Exception as member_err:
                                combined_parts.append(
                                    f"XAF Member: {member_name}\n"
                                    f"Failed to parse member: {member_err}"
                                )
                        return "\n\n".join(combined_parts)
                except Exception as zip_err:
                    logger.error(f"Error reading zipped XAF {file_path.name}: {zip_err}")
                    return f"Error reading zipped XAF: {zip_err}"

            logger.error(f"Error parsing XAF XML: {file_path.name}")
            return f"Error reading XAF: invalid XML in {file_path.name}"
        except Exception as e:
            logger.error(f"Error extracting text from XAF: {e}")
            return f"Error reading XAF: {str(e)}"

    def extract_text_from_html(self, file_path: Path) -> str:
        """
        Extract text content from HTML file.

        Args:
            file_path: Path to HTML file

        Returns:
            str: Plain text extracted from HTML
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f, 'html.parser')

                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()

                # Get text
                text = soup.get_text(separator='\n', strip=True)

                # Clean up excessive whitespace
                lines = (line.strip() for line in text.splitlines())
                text = '\n'.join(line for line in lines if line)

                return text

        except Exception as e:
            logger.error(f"Error extracting text from HTML: {e}")
            return f"Error reading HTML: {str(e)}"

    def extract_text_from_image(self, file_path: Path) -> str:
        """
        Extract text from image using OCR (Optical Character Recognition).

        Args:
            file_path: Path to image file

        Returns:
            str: Extracted text from image
        """
        try:
            # Try to use pytesseract for OCR
            try:
                from PIL import Image
                import pytesseract

                # Open image
                image = Image.open(file_path)

                # Perform OCR
                text = pytesseract.image_to_string(image)

                if text.strip():
                    logger.info(f"Extracted {len(text)} characters from image {file_path.name}")
                    return f"Image: {file_path.name}\nExtracted text:\n\n{text}"
                else:
                    return f"Image file: {file_path.name} (no text detected)\nType: {file_path.suffix}\nSize: {file_path.stat().st_size} bytes"

            except ImportError as e:
                logger.warning(f"OCR libraries not installed: {e}")
                return f"Image file: {file_path.name} (OCR not available - install pytesseract and Pillow)\nType: {file_path.suffix}\nSize: {file_path.stat().st_size} bytes"
            except Exception as e:
                logger.warning(f"OCR failed for {file_path.name}: {e}")
                return f"Image file: {file_path.name} (OCR failed)\nType: {file_path.suffix}\nSize: {file_path.stat().st_size} bytes"

        except Exception as e:
            logger.error(f"Error processing image: {e}")
            return f"Error reading image: {str(e)}"

    def extract_text_from_docx(self, file_path: Path) -> str:
        """
        Extract text from DOCX file.

        Args:
            file_path: Path to DOCX file

        Returns:
            str: Extracted text content
        """
        try:
            # Try to use python-docx if available
            try:
                import docx
                doc = docx.Document(file_path)
                text_parts = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)
                return '\n\n'.join(text_parts)
            except ImportError:
                logger.warning("python-docx not installed, using basic DOCX parsing")
                # Fallback: DOCX is a ZIP file containing XML
                import zipfile
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    try:
                        xml_content = zip_ref.read('word/document.xml')
                        soup = BeautifulSoup(xml_content, 'xml')
                        paragraphs = soup.find_all('t')
                        return '\n'.join(p.get_text() for p in paragraphs)
                    except Exception as e:
                        return f"DOCX file (text extraction requires python-docx): {file_path.name}"
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
            return f"Error reading DOCX: {str(e)}"

    def extract_text_from_xlsx(self, file_path: Path) -> str:
        """
        Extract text from XLSX file.

        Args:
            file_path: Path to XLSX file

        Returns:
            str: Extracted text content
        """
        try:
            suffix = file_path.suffix.lower()
            if suffix == ".xls":
                try:
                    import xlrd

                    book = xlrd.open_workbook(file_path)
                    text_parts = []
                    for sheet in book.sheets():
                        text_parts.append(f"Sheet: {sheet.name}\n")
                        for row_idx in range(sheet.nrows):
                            row_values = sheet.row_values(row_idx)
                            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row_values)
                            if row_text.strip():
                                text_parts.append(row_text)
                        text_parts.append("")
                    return "\n".join(text_parts)
                except ImportError:
                    return (
                        f"XLS file (text extraction requires xlrd): {file_path.name}\n"
                        "Install with: pip install xlrd"
                    )

            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text_parts = []

                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    text_parts.append(f"Sheet: {sheet_name}\n")

                    for row in sheet.iter_rows(values_only=True):
                        row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                        if row_text.strip():
                            text_parts.append(row_text)
                    text_parts.append('')

                return '\n'.join(text_parts)
            except ImportError:
                logger.warning("openpyxl not installed")
                return f"XLSX file (text extraction requires openpyxl): {file_path.name}\nInstall with: pip install openpyxl"
        except Exception as e:
            logger.error(f"Error extracting text from Excel: {e}")
            return f"Error reading Excel: {str(e)}"

    def extract_text_from_pptx(self, file_path: Path) -> str:
        """
        Extract text from PPTX file.

        Args:
            file_path: Path to PPTX file

        Returns:
            str: Extracted text content
        """
        try:
            # Try to use python-pptx if available
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []

                for slide_num, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"Slide {slide_num}:")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                    text_parts.append('')

                return '\n'.join(text_parts)
            except ImportError:
                logger.warning("python-pptx not installed, using basic PPTX parsing")
                return f"PPTX file (text extraction requires python-pptx): {file_path.name}\nInstall with: pip install python-pptx"
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
            return f"Error reading PPTX: {str(e)}"

    def extract_text_from_pdf(self, file_path: Path) -> str:
        """
        Extract text from PDF file.

        Args:
            file_path: Path to PDF file

        Returns:
            str: Extracted text content
        """
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            text_parts = []

            total_pages = len(reader.pages)
            logger.info(f"Extracting {total_pages} pages from {file_path.name}")

            # Extract text from each page
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text() or ""
                    if not page_text.strip():
                        page_text = self._ocr_text_from_pdf_page_images(page, file_path.name, page_num)
                    if not page_text.strip():
                        page_text = self._ocr_text_from_pdf_page_render(file_path, page_num)
                    if page_text.strip():
                        text_parts.append(f"--- Page {page_num} ---\n{page_text}\n")
                    # Lightweight page counter for long PDFs
                    if total_pages > 20 and (page_num == 1 or page_num % 250 == 0 or page_num == total_pages):
                        print(f"    {file_path.name}: page {page_num}/{total_pages}")
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num} of {file_path.name}: {e}")
                    text_parts.append(f"--- Page {page_num} ---\n[Text extraction failed]\n")

            if not text_parts:
                return f"PDF file: {file_path.name} (no text could be extracted)"

            full_text = '\n'.join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from {len(reader.pages)} pages of {file_path.name}")
            return full_text

        except ImportError:
            logger.error("pypdf not installed. Install with: pip install pypdf")
            return f"PDF file: {file_path.name} (pypdf library not available)\nInstall with: pip install pypdf"
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            return f"Error reading PDF: {str(e)}"

    def extract_pages_from_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text from PDF file, page by page.

        Args:
            file_path: Path to PDF file

        Returns:
            List[Dict]: List of pages [{'page': 1, 'content': '...'}, ...]
        """
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            pages = []

            total_pages = len(reader.pages)
            logger.info(f"Extracting {total_pages} pages from {file_path.name}")
            print(f"    Started extraction: {file_path.name} ({total_pages} pages)")

            for page_num, page in enumerate(reader.pages, 1):
                if page_num % 50 == 0:
                    print(f"    Processing {file_path.name}: Page {page_num}/{total_pages}...")
                
                try:
                    page_text = page.extract_text() or ""
                    if not page_text.strip():
                        page_text = self._ocr_text_from_pdf_page_images(page, file_path.name, page_num)
                    if not page_text.strip():
                        page_text = self._ocr_text_from_pdf_page_render(file_path, page_num)
                    if page_text and page_text.strip():
                        pages.append({
                            "page": page_num,
                            "content": page_text
                        })
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num} of {file_path.name}: {e}")

            return pages

        except ImportError:
            logger.error("pypdf not installed")
            return []
        except Exception as e:
            logger.error(f"Error extracting pages from PDF: {e}")
            return []

    def _ocr_text_from_pdf_page_images(self, page: Any, source_filename: str, page_num: int) -> str:
        """
        OCR fallback for scanned PDF pages by extracting page images.
        """
        try:
            import pytesseract
            from PIL import Image
        except Exception:
            return ""

        texts: List[str] = []
        try:
            page_images = list(page.images)
        except Exception as e:
            logger.debug(f"No usable page.images for OCR ({source_filename} p{page_num}): {e}")
            return ""

        for image_obj in page_images:
            try:
                if hasattr(image_obj, "image"):
                    pil_image = image_obj.image
                else:
                    pil_image = Image.open(io.BytesIO(image_obj.data))

                ocr_text = pytesseract.image_to_string(pil_image) or ""
                if ocr_text.strip():
                    texts.append(ocr_text.strip())
            except Exception as e:
                logger.debug(f"OCR failed for image on {source_filename} p{page_num}: {e}")

        if texts:
            logger.info(
                f"OCR fallback used for {source_filename} page {page_num}: "
                f"{sum(len(t) for t in texts)} chars"
            )
        return "\n".join(texts)

    def _ocr_text_from_pdf_page_render(self, file_path: Path, page_num: int) -> str:
        """
        OCR fallback by rendering PDF page to bitmap using pypdfium2.
        """
        try:
            import pytesseract
            import pypdfium2 as pdfium
        except Exception:
            return ""

        try:
            pdf = pdfium.PdfDocument(str(file_path))
            page_index = max(0, page_num - 1)
            if page_index >= len(pdf):
                return ""

            page = pdf[page_index]
            # scale=2 gives readable OCR quality while keeping it fast
            bitmap = page.render(scale=2)
            pil_image = bitmap.to_pil()
            text = (pytesseract.image_to_string(pil_image) or "").strip()
            if text:
                logger.info(
                    f"OCR render fallback used for {file_path.name} page {page_num}: {len(text)} chars"
                )
            return text
        except Exception as e:
            logger.debug(f"OCR render fallback failed for {file_path.name} p{page_num}: {e}")
            return ""

    def extract_text_from_json(self, file_path: Path) -> str:
        """
        Extract text from JSON file.

        Args:
            file_path: Path to JSON file

        Returns:
            str: Formatted JSON content
        """
        try:
            import json
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)
                # Pretty print JSON for readability
                return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            logger.error(f"Error reading JSON: {e}")
            return f"Error reading JSON: {str(e)}"

    def extract_text_from_file(self, file_path: Path) -> str:
        """
        Extract text from any supported file format.

        Args:
            file_path: Path to file

        Returns:
            str: Extracted text content
        """
        file_type = self.detect_file_type(file_path)

        logger.info(f"Extracting text from {file_type} file: {file_path.name}")

        if file_type == 'pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_type == 'csv':
            return self.extract_text_from_csv(file_path)
        elif file_type == 'xaf':
            return self.extract_text_from_xaf(file_path)
        elif file_type == 'xml':
            return self.extract_text_from_xml(file_path)
        elif file_type == 'html':
            return self.extract_text_from_html(file_path)
        elif file_type == 'image':
            return self.extract_text_from_image(file_path)
        elif file_type == 'docx':
            return self.extract_text_from_docx(file_path)
        elif file_type == 'xlsx' or file_type == 'xls':
            return self.extract_text_from_xlsx(file_path)
        elif file_type == 'pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_type == 'json':
            return self.extract_text_from_json(file_path)
        elif file_type == 'text':
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception as e:
                return f"Error reading text file: {str(e)}"
        elif file_type == 'unknown':
            # For unknown file types (like Confluence attachments without extensions),
            # try to detect content type by reading the file
            try:
                # First, try reading as text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(10000)  # Read first 10KB

                    # Check if it looks like JSON
                    if content.strip().startswith('{') or content.strip().startswith('['):
                        # Read full file as JSON/text
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f2:
                            return f2.read()

                    # Check if it looks like HTML
                    elif '<html' in content.lower() or '<!doctype html' in content.lower():
                        return self.extract_text_from_html(file_path)

                    # Check if it looks like XML
                    elif content.strip().startswith('<?xml') or content.strip().startswith('<'):
                        return self.extract_text_from_xml(file_path)

                    # Otherwise treat as plain text
                    else:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f2:
                            text = f2.read()
                            # If file is mostly readable text, return it
                            if text.strip():
                                return text
                            else:
                                return f"Empty or binary file: {file_path.name}"

            except Exception as e:
                logger.debug(f"Could not read unknown file {file_path.name} as text: {e}")
                # File is likely binary, skip it
                return f"Binary file (skipped): {file_path.name}"
        else:
            return f"Unsupported file type: {file_type}"

    def _render_xml_tree(self, root: ET.Element, header: str) -> str:
        """Convert XML tree into line-based text for indexing."""
        text_parts = [header, ""]
        for elem in root.iter():
            tag_name = elem.tag.split('}')[-1]
            text_value = (elem.text or "").strip()
            if text_value:
                text_parts.append(f"{tag_name}: {text_value}")
            if elem.attrib:
                for key, value in elem.attrib.items():
                    text_parts.append(f"{tag_name}.@{key}: {value}")
        return "\n".join(text_parts)

    def generate_summary(self, content: str, file_name: str) -> str:
        """
        Generate a summary of document content.
        Falls back to a short excerpt if AI is not available.

        Args:
            content: Text content to summarize
            file_name: Name of source file

        Returns:
            str: AI-generated summary or simple excerpt
        """
        if not settings.anthropic_api_key or settings.anthropic_api_key.startswith('dummy'):
            excerpt = content[:500] + "..." if len(content) > 500 else content
            logger.debug(f"Using excerpt for {file_name} (no AI key)")
            return excerpt

        try:
            max_content_length = 100000
            if len(content) > max_content_length:
                content = content[:max_content_length] + "... [truncated]"

            prompt = f"""You are analyzing a document named "{file_name}".

Please create a short searchable summary of this document that:
1. Captures the main topics and key information
2. Includes important technical terms, names, and identifiers
3. Describes the document's purpose and content
4. Uses clear language

Document content:
{content}

Provide a summary (2-4 paragraphs) that helps someone find this document in search."""

            message = self.anthropic_client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=1000,
                messages=[{"role": "user", "content": prompt}]
            )

            summary = message.content[0].text
            logger.info(f"Generated AI summary for {file_name} ({len(summary)} chars)")
            return summary

        except Exception as e:
            logger.warning(f"AI summary failed for {file_name}, using excerpt: {e}")
            return content[:500] + "..." if len(content) > 500 else content

    def process_document(
        self,
        file_path: Path,
        document_id: Optional[str] = None,
        category: Optional[DocumentCategory] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Process a document of any format.

        Args:
            file_path: Path to document
            document_id: Optional document ID (generated if not provided)
            category: Document category
            metadata: Additional metadata

        Returns:
            dict: Processing result with document data
        """
        if document_id is None:
            document_id = str(uuid.uuid4())

        if category is None:
            category = DocumentCategory.OTHER

        if metadata is None:
            metadata = {}

        logger.info(f"Processing document: {file_path.name} (ID: {document_id})")

        result = {
            "document_id": document_id,
            "filename": file_path.name,
            "file_type": self.detect_file_type(file_path),
            "status": ProcessingStatus.UPLOADED,
            "category": category.value if hasattr(category, 'value') else str(category),
            "metadata": metadata,
            "content": None,
            "summary": None,
            "error_message": None,
            "processed_at": datetime.utcnow().isoformat()
        }

        try:
            # Extract text content
            result["status"] = ProcessingStatus.PARSING
            content = self.extract_text_from_file(file_path)
            result["content"] = content

            # Generate summary
            result["status"] = ProcessingStatus.SUMMARIZING
            summary = self.generate_summary(content, file_path.name)
            result["summary"] = summary

            # Mark as ready
            result["status"] = ProcessingStatus.READY

            logger.info(f"Successfully processed {file_path.name}")
            return result

        except Exception as e:
            logger.error(f"Error processing document {file_path.name}: {e}")
            result["status"] = ProcessingStatus.ERROR
            result["error_message"] = str(e)
            return result


# Singleton instance
_processor: Optional[MultiFormatProcessor] = None


def get_multi_format_processor() -> MultiFormatProcessor:
    """
    Get the global MultiFormatProcessor instance.

    Returns:
        MultiFormatProcessor: The processor instance
    """
    global _processor
    if _processor is None:
        _processor = MultiFormatProcessor()
    return _processor
